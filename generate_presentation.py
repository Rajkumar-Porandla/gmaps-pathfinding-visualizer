from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    # Color Palette: G-Maps Theme
    BLUE = RGBColor(26, 115, 232)      # #1A73E8
    DARK_BLUE = RGBColor(21, 87, 176)  # #1557B0
    SLATE = RGBColor(95, 99, 104)      # #5F6368
    CHARCOAL = RGBColor(32, 33, 36)    # #202124
    LIGHT_BG = RGBColor(248, 249, 250) # #F8F9FA
    WHITE = RGBColor(255, 255, 255)

    # Helper: Set slide background color
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper: Add slide title
    def add_slide_title(slide, text):
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Arial'
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = CHARCOAL
        return title_box

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, LIGHT_BG)

    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
    tf = title_box.text_frame
    
    p = tf.paragraphs[0]
    p.text = "GOOGLE MAPS INSPIRED ROUTE VISUALIZER"
    p.font.name = 'Arial'
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = "A High-Performance Pathfinding and Routing Simulation Platform"
    p2.font.name = 'Arial'
    p2.font.size = Pt(18)
    p2.font.italic = True
    p2.font.color.rgb = SLATE
    
    p3 = tf.add_paragraph()
    p3.text = "\nSkill Development Project | B.Tech CSE 3rd Year\nSubmitted By: Porandla Rajkumar"
    p3.font.name = 'Arial'
    p3.font.size = Pt(14)
    p3.font.bold = True
    p3.font.color.rgb = CHARCOAL

    # Slide 2: Objectives & Context
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)
    add_slide_title(slide, "Project Context & Key Objectives")
    
    tb = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    bullets = [
        ("Core Scope", "Developed as a Skill Development Project to bridge the gap between theoretical graph algorithms and real-time visualization."),
        ("Google Maps Theming", "Styled with real-world geographical design cues: location pins, construction roadblocks, and traffic delay markers."),
        ("Navigation Logic", "Simulates unweighted (BFS, DFS) and weighted (Dijkstra, A* Search) algorithms in dynamic conditions."),
        ("Performance-First", "Enforces layout isolation (CSS Containment) and Direct-DOM painting to eliminate React reconciliation lag during high-frequency visual updates.")
    ]
    
    for title, desc in bullets:
        p = tf.add_paragraph()
        p.space_after = Pt(14)
        run_title = p.add_run()
        run_title.text = f"•  {title}: "
        run_title.font.bold = True
        run_title.font.size = Pt(16)
        run_title.font.color.rgb = BLUE
        
        run_desc = p.add_run()
        run_desc.text = desc
        run_desc.font.size = Pt(15)
        run_desc.font.color.rgb = CHARCOAL

    # Slide 3: Technology Stack
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)
    add_slide_title(slide, "Technology Stack & Rationale")
    
    tb = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    techs = [
        ("React (v19)", "Component reuse, state hooks for controls, and modular layout encapsulation."),
        ("TypeScript", "Enforces static type checking, resolving verbatimModuleSyntax imports and null references."),
        ("TailwindCSS v4", "Utilizes modern CSS variables for light/dark skins and responsive flex/grid wrappers."),
        ("Vite (v6)", "Provides instantaneous HMR and extremely lightweight, fast ESM production packaging."),
        ("Framer Motion", "Drives physics-based spring animations for overlay guides and transitions.")
    ]
    
    for tech, reason in techs:
        p = tf.add_paragraph()
        p.space_after = Pt(12)
        r = p.add_run()
        r.text = f"•  {tech}: "
        r.font.bold = True
        r.font.size = Pt(16)
        r.font.color.rgb = DARK_BLUE
        
        r2 = p.add_run()
        r2.text = reason
        r2.font.size = Pt(15)
        r2.font.color.rgb = CHARCOAL

    # Slide 4: Pathfinding Algorithms (BFS vs DFS)
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)
    add_slide_title(slide, "Pathfinding: BFS & DFS (Unweighted)")
    
    tb = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.add_paragraph()
    p1.text = "Breadth-First Search (BFS)"
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = BLUE
    
    p2 = tf.add_paragraph()
    p2.text = "  - Strategy: Explores level-by-level using a FIFO Queue.\n  - Shortest Path: Guarantees shortest path on unweighted grids.\n  - Complexity: O(V + E) Time | O(V) Space."
    p2.font.size = Pt(14)
    p2.font.color.rgb = CHARCOAL
    p2.space_after = Pt(18)
    
    p3 = tf.add_paragraph()
    p3.text = "Depth-First Search (DFS)"
    p3.font.bold = True
    p3.font.size = Pt(18)
    p3.font.color.rgb = BLUE
    
    p4 = tf.add_paragraph()
    p4.text = "  - Strategy: Explores as far as possible down each branch before backtracking. Uses a LIFO Stack.\n  - Shortest Path: Does NOT guarantee shortest path.\n  - Complexity: O(V + E) Time | O(V) Space."
    p4.font.size = Pt(14)
    p4.font.color.rgb = CHARCOAL

    # Slide 5: Pathfinding Algorithms (Dijkstra vs A*)
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)
    add_slide_title(slide, "Pathfinding: Dijkstra & A* Search (Weighted)")
    
    tb = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.add_paragraph()
    p1.text = "Dijkstra's Algorithm"
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = BLUE
    
    p2 = tf.add_paragraph()
    p2.text = "  - Strategy: Extracts lowest-cost unvisited nodes using a binary Min-Heap priority queue.\n  - Shortest Path: Guarantees shortest path. Adapts to heavy traffic congestion zones.\n  - Complexity: O((V + E) log V) Time | O(V) Space."
    p2.font.size = Pt(14)
    p2.font.color.rgb = CHARCOAL
    p2.space_after = Pt(18)
    
    p3 = tf.add_paragraph()
    p3.text = "A* Search Algorithm"
    p3.font.bold = True
    p3.font.size = Pt(18)
    p3.font.color.rgb = BLUE
    
    p4 = tf.add_paragraph()
    p4.text = "  - Strategy: Minimizes f(n) = g(n) + h(n), guided by Manhattan distance heuristic.\n  - Shortest Path: Guarantees shortest path. Explores highly targeted areas, beating Dijkstra.\n  - Complexity: O(E log V) Average Time | O(b^d) Space."
    p4.font.size = Pt(14)
    p4.font.color.rgb = CHARCOAL

    # Slide 6: Maze and Obstacle Generation
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)
    add_slide_title(slide, "Maze & Terrain Generation")
    
    tb = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    mazes = [
        ("Recursive Division", "Fractally divides the grid with horizontal/vertical walls. Places walls on even indices and passages on odd indices, forming consistent streets."),
        ("Random Roadblocks", "Scatters solid barrier roadblock walls across open grid paths based on a 30% density factor."),
        ("Random Traffic Jams", "Distributes heavy congestion weights (5x travel cost multiplier) across 25% of the grid to evaluate weighted path routing.")
    ]
    
    for title, desc in mazes:
        p = tf.add_paragraph()
        p.space_after = Pt(14)
        r = p.add_run()
        r.text = f"•  {title}: "
        r.font.bold = True
        r.font.size = Pt(16)
        r.font.color.rgb = DARK_BLUE
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(15)
        r2.font.color.rgb = CHARCOAL

    # Slide 7: Playback Controller & Features
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)
    add_slide_title(slide, "Playback Controller & Core Features")
    
    tb = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    feats = [
        ("Step-by-Step Controller", "Play, Pause, Resume, Stop, and Restart routing animations. Suspends operations dynamically using index refs."),
        ("Speed and Grid Customizers", "6 speed presets (Very Slow to Instant) and 6 grid options (20x20 up to 60x40) with auto-resizing."),
        ("Comparison Benchmarking", "Compares all pathfinders concurrently, reporting CPU speeds, visited nodes, path cost, and RAM. Exports to JSON/CSV."),
        ("JSON Save & Load", "Encodes wall/weight coordinates and dimensions to JSON files, allowing layouts sharing.")
    ]
    
    for title, desc in feats:
        p = tf.add_paragraph()
        p.space_after = Pt(12)
        r = p.add_run()
        r.text = f"•  {title}: "
        r.font.bold = True
        r.font.size = Pt(16)
        r.font.color.rgb = BLUE
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(15)
        r2.font.color.rgb = CHARCOAL

    # Slide 8: Viva Inspector & Demonstration Tools
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)
    add_slide_title(slide, "Viva Mode & Demonstration Panel")
    
    tb = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "Designed specifically to aid students during academic project evaluations:"
    p.font.italic = True
    p.font.size = Pt(15)
    p.font.color.rgb = SLATE
    p.space_after = Pt(14)
    
    viva_points = [
        ("Live Stack/Queue View", "Renders index contents of the active Queue (BFS) or Stack (DFS) at each evaluation step."),
        ("Heap Priority Queue Inspector", "Displays node weights and f/g/h values sorted in the binary Min-Heap during Dijkstra/A* execution."),
        ("Interactive Suspend", "Evaluators can pause animations to inspect exact parent-child relations and distance variables of current nodes.")
    ]
    
    for title, desc in viva_points:
        p = tf.add_paragraph()
        p.space_after = Pt(12)
        r = p.add_run()
        r.text = f"•  {title}: "
        r.font.bold = True
        r.font.size = Pt(16)
        r.font.color.rgb = DARK_BLUE
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(15)
        r2.font.color.rgb = CHARCOAL

    # Slide 9: Performance Optimizations
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, WHITE)
    add_slide_title(slide, "Performance & Memory Tuning")
    
    tb = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    opts = [
        ("Direct-DOM Paint Injections", "Bypasses React virtual DOM reconciliation during playback. Paints classes directly onto cells, avoiding browser reflow lag and maintaining a steady 60 FPS on large 60x40 grids."),
        ("CSS Layout Containment", "Applies 'contain: layout style paint' rules on nodes, isolating style recalcs to modified grid coordinates."),
        ("Optimized Min-Heap", "Replaced standard array sort queries with binary heaps, dropping Dijkstra runtime from O(V^2) to O((V+E) log V)."),
        ("React.memo & useCallback Hooks", "Memoizes Node components to block redundant re-renders of unmodified cells.")
    ]
    
    for title, desc in opts:
        p = tf.add_paragraph()
        p.space_after = Pt(12)
        r = p.add_run()
        r.text = f"•  {title}: "
        r.font.bold = True
        r.font.size = Pt(16)
        r.font.color.rgb = BLUE
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(15)
        r2.font.color.rgb = CHARCOAL

    # Slide 10: Conclusion
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, LIGHT_BG)
    
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
    tf = title_box.text_frame
    
    p = tf.paragraphs[0]
    p.text = "THANK YOU!"
    p.font.name = 'Arial'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = "Google Maps Inspired Path Finding Visualizer\nQuestions & Demonstration"
    p2.font.name = 'Arial'
    p2.font.size = Pt(18)
    p2.font.color.rgb = CHARCOAL
    
    prs.save("G-Maps_Route_Visualizer_Presentation.pptx")
    print("Presentation generated successfully as G-Maps_Route_Visualizer_Presentation.pptx")

if __name__ == '__main__':
    create_presentation()
